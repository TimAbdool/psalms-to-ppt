import os
import re
from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from docx import Document

app = Flask(__name__)

def extract_text_from_docx(filepath):
    doc = Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_verses(text):
    verses = []
    current_psalm = ""

    lines = text.splitlines()

    for line in lines:
        line = line.strip()

        psalm_match = re.match(r'Psalm\s+(\d+)', line, re.IGNORECASE)
        if psalm_match:
            current_psalm = f"Psalm {psalm_match.group(1)}"
            continue

        verse_match = re.match(r'^(\d+)\s+(.*)', line)
        if verse_match:
            verse_number = verse_match.group(1)
            verse_text = verse_match.group(2)
            full_verse = f"{current_psalm} {verse_number}\n{verse_text}"
            verses.append(full_verse)

    return verses

def auto_fit_text(paragraph, text):
    max_font_size = 72
    min_font_size = 36
    paragraph.text = text
    paragraph.font.name = "Times New Roman"
    paragraph.alignment = PP_ALIGN.CENTER

    for size in range(max_font_size, min_font_size - 1, -2):
        paragraph.font.size = Pt(size)
        if len(text) < 200:
            break

def create_powerpoint(verses, output_path):
    prs = Presentation()

    for verse in verses:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        textbox = slide.shapes.add_textbox(Pt(50), Pt(150), Pt(860), Pt(400))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()

        p = text_frame.paragraphs[0]
        auto_fit_text(p, verse)

    prs.save(output_path)

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        file = request.files["file"]
        filepath = file.filename
        file.save(filepath)

        if filepath.endswith(".docx"):
            text = extract_text_from_docx(filepath)
        else:
            text = file.read().decode("utf-8")

        verses = extract_verses(text)
        output_file = "Psalms_Output.pptx"
        create_powerpoint(verses, output_file)

        return send_file(output_file, as_attachment=True)

    return render_template("index.html")

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port)
